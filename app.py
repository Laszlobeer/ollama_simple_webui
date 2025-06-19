from flask import Flask, render_template, request, jsonify, send_from_directory
import os
import json
import subprocess
import requests
from duckduckgo_search import DDGS
from PyPDF2 import PdfReader
from docx import Document
import re
from urllib.parse import urlparse
import html
import pandas as pd
from pptx import Presentation
import pytesseract
from PIL import Image
import speech_recognition as sr
from io import BytesIO
import openpyxl

app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = 'uploads'

# Theme loader
@app.context_processor
def inject_theme():
    return dict(theme=request.cookies.get('theme', 'dark'))

# Scan Ollama models
@app.route('/api/models')
def get_models():
    try:
        # First try with JSON output
        result = subprocess.run(['ollama', 'list', '--json'], capture_output=True, text=True)
        if result.returncode == 0 and result.stdout.strip():
            models = [model['name'] for model in json.loads(result.stdout)]
            return jsonify(models)
        
        # Fallback to table parsing if JSON fails
        result = subprocess.run(['ollama', 'list'], capture_output=True, text=True)
        if result.returncode == 0 and result.stdout.strip():
            models = []
            for line in result.stdout.splitlines()[1:]:
                if line.strip():
                    model_name = line.split()[0]
                    if '/' in model_name:  # Handle names like 'library/llama3'
                        model_name = model_name.split('/')[-1]
                    models.append(model_name)
            return jsonify(models)
        
        return jsonify([])  # Return empty list if no models found
    except Exception as e:
        return jsonify({'error': str(e)}), 500

# Process file content
def extract_text_from_file(filepath):
    text = ""
    try:
        # Get file extension
        _, ext = os.path.splitext(filepath)
        ext = ext.lower()
        
        # Handle different file types
        if ext == '.pdf':
            with open(filepath, 'rb') as f:
                reader = PdfReader(f)
                for page in reader.pages:
                    text += page.extract_text() + "\n"
                    
        elif ext in ['.docx', '.doc']:
            doc = Document(filepath)
            for para in doc.paragraphs:
                text += para.text + "\n"
                
        elif ext == '.txt':
            with open(filepath, 'r', encoding='utf-8') as f:
                text = f.read()
                
        elif ext in ['.xlsx', '.xls']:
            # Handle Excel files
            if ext == '.xlsx':
                wb = openpyxl.load_workbook(filepath)
                for sheet_name in wb.sheetnames:
                    sheet = wb[sheet_name]
                    text += f"--- Sheet: {sheet_name} ---\n"
                    for row in sheet.iter_rows(values_only=True):
                        text += "\t".join([str(cell) if cell is not None else "" for cell in row]) + "\n"
                    text += "\n"
            else:  # .xls
                df = pd.read_excel(filepath, sheet_name=None)
                for sheet_name, sheet_data in df.items():
                    text += f"--- Sheet: {sheet_name} ---\n"
                    text += sheet_data.to_string(index=False) + "\n\n"
                    
        elif ext == '.csv':
            df = pd.read_csv(filepath)
            text = df.to_string(index=False)
            
        elif ext in ['.pptx', '.ppt']:
            prs = Presentation(filepath)
            for i, slide in enumerate(prs.slides):
                text += f"--- Slide {i+1} ---\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
                text += "\n"
                
        elif ext in ['.png', '.jpg', '.jpeg', '.bmp', '.tiff']:
            # OCR for images
            img = Image.open(filepath)
            text = pytesseract.image_to_string(img)
            
        elif ext in ['.py', '.js', '.html', '.css', '.java', '.cpp', '.c', '.php', '.rb', '.go', '.rs', '.ts', '.sh']:
            # Code files
            with open(filepath, 'r', encoding='utf-8') as f:
                text = f.read()
                
        elif ext in ['.wav', '.mp3', '.flac']:
            # Audio transcription
            r = sr.Recognizer()
            with sr.AudioFile(filepath) as source:
                audio = r.record(source)
            text = r.recognize_google(audio)
            
        else:
            # Try to read as text
            try:
                with open(filepath, 'r', encoding='utf-8') as f:
                    text = f.read()
            except:
                text = f"[Unsupported file format: {ext}]"
        
        # Clean and truncate text
        text = re.sub(r'\s+', ' ', text).strip()
        return text[:15000]  # Truncate to 15k characters
    except Exception as e:
        return f"[Error processing file: {str(e)}]"

# Chat endpoint
@app.route('/api/chat', methods=['POST'])
def chat():
    try:
        data = request.json
        model = data.get('model', 'llama3')
        messages = data.get('messages', [])
        options = data.get('options', {})
        web_search = data.get('webSearch', False)
        
        # Ollama API integration
        ollama_url = 'http://localhost:11434/api/chat'
        payload = {
            "model": model,
            "messages": messages,
            "stream": False,
            "options": {
                "temperature": options.get('temperature', 0.8),
                "top_p": options.get('topP', 0.9),
                "top_k": options.get('topK', 40),
                "repeat_penalty": options.get('repeat_penalty', 1.1),
                "num_predict": options.get('maxTokens', 2048),
                "num_ctx": options.get('contextLength', 2048)
            }
        }
        
        # Add system prompt if provided
        if 'system' in data:
            payload['system'] = data['system']
            
        # Add web search functionality if enabled
        web_search_results = []
        webSearchNote = ""
        if web_search:
            # Find the most recent user message
            user_messages = [msg['content'] for msg in messages if msg['role'] == 'user']
            if user_messages:
                last_user_message = user_messages[-1]
                
                # Perform web search
                try:
                    with DDGS() as ddgs:
                        results = ddgs.text(last_user_message, max_results=5)
                        web_search_results = [
                            {
                                "title": r["title"],
                                "url": r["href"],
                                "snippet": r["body"],
                                "domain": urlparse(r["href"]).netloc
                            } for r in results
                        ]
                except Exception as e:
                    app.logger.error(f"Web search error: {str(e)}")
                
                # Format search results for AI context
                if web_search_results:
                    webSearchNote = "I found these sources to help answer your question:"
                    search_context = "### Current Web Search Results:\n"
                    search_context += "Use these search results to provide a direct answer to the user's question. " \
                                     "Cite sources using their domain names in parentheses. " \
                                     "When possible, provide specific facts, figures, or quotes from the sources.\n\n"
                    
                    for i, result in enumerate(web_search_results):
                        search_context += f"{i+1}. [{result['title']}]({result['url']})\n"
                        search_context += f"   Summary: {result['snippet']}\n\n"
                    
                    # Add search context to the payload
                    if 'system' in payload:
                        payload['system'] += "\n\n" + search_context
                    else:
                        payload['system'] = search_context
            else:
                webSearchNote = "No relevant sources found for your question"
        
        # Process file content if referenced in messages
        for msg in messages:
            if msg['role'] == 'user' and '📄 Attached file' in msg['content']:
                # Extract filename from message
                match = re.search(r'📄 Attached file: (.+?)\n', msg['content'])
                if match:
                    filename = match.group(1)
                    filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
                    if os.path.exists(filepath):
                        file_content = extract_text_from_file(filepath)
                        msg['content'] += f"\n\nFile content:\n{file_content}"
        
        response = requests.post(ollama_url, json=payload)
        response.raise_for_status()
        
        # Add web search results to response
        response_data = response.json()
        if web_search_results:
            response_data['web_search_results'] = web_search_results
            response_data['webSearchNote'] = webSearchNote
            
        return jsonify(response_data)
    except Exception as e:
        return jsonify({'error': str(e)}), 500

# System prompt handler
@app.route('/api/prompt', methods=['POST'])
def update_prompt():
    try:
        new_prompt = request.json.get('prompt')
        # Save to file
        with open('system_prompt.txt', 'w') as f:
            f.write(new_prompt)
        return jsonify({'status': 'success'})
    except Exception as e:
        return jsonify({'error': str(e)}), 500

# Get current system prompt
@app.route('/api/prompt', methods=['GET'])
def get_prompt():
    try:
        if os.path.exists('system_prompt.txt'):
            with open('system_prompt.txt', 'r') as f:
                prompt = f.read()
            return jsonify({'prompt': prompt})
        return jsonify({'prompt': 'You are a helpful AI assistant.'})
    except Exception as e:
        return jsonify({'error': str(e)}), 500

# File upload handler
@app.route('/upload', methods=['POST'])
def upload_file():
    try:
        if 'file' not in request.files:
            return jsonify({'error': 'No file part'}), 400
            
        file = request.files['file']
        if file.filename == '':
            return jsonify({'error': 'No selected file'}), 400
            
        os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], file.filename)
        file.save(file_path)
        
        # Extract text content for AI
        file_content = extract_text_from_file(file_path)
        
        return jsonify({
            'filename': file.filename,
            'path': file_path,
            'content': file_content
        })
    except Exception as e:
        return jsonify({'error': str(e)}), 500

# Serve uploaded files
@app.route('/uploads/<filename>')
def uploaded_file(filename):
    return send_from_directory(app.config['UPLOAD_FOLDER'], filename)

# History management
@app.route('/api/history', methods=['DELETE'])
def clear_history():
    try:
        open('history.json', 'w').close()
        return jsonify({'status': 'cleared'})
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/')
def index():
    theme = request.cookies.get('theme', 'dark')
    return render_template('index.html', theme=theme)

if __name__ == '__main__':
    app.run(debug=True, port=5001)